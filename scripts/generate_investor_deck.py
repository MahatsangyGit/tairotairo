#!/usr/bin/env python3
"""Tairo ampio — Investor deck (15 slides, 16:9)."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData
import os

# ── Palette ──────────────────────────────────────────────────────────────────
TEAL = RGBColor(0x06, 0x94, 0x94)
PEACH = RGBColor(0xFA, 0xD0, 0xC9)
INK = RGBColor(0x0F, 0x17, 0x2A)
SLATE = RGBColor(0x64, 0x74, 0x8B)
BG = RGBColor(0xF8, 0xFA, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEAL_DARK = RGBColor(0x04, 0x6E, 0x6E)
TEAL_SOFT = RGBColor(0xE6, 0xF6, 0xF6)
PEACH_SOFT = RGBColor(0xFE, 0xF3, 0xF0)
RED_SOFT = RGBColor(0xFE, 0xF2, 0xF2)
GREEN_SOFT = RGBColor(0xEC, 0xFD, 0xF5)
AMBER_SOFT = RGBColor(0xFF, 0xFB, 0xEB)
RED = RGBColor(0xDC, 0x26, 0x26)
GREEN = RGBColor(0x05, 0x96, 0x69)
AMBER = RGBColor(0xD9, 0x77, 0x06)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
TOTAL = 15


def set_run(run, size=18, bold=False, color=INK, font="Calibri"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_text(shape, text, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT, font="Calibri"):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color, font=font)
    return tf


def add_para(tf, text, size=16, bold=False, color=INK, align=PP_ALIGN.LEFT, space_before=6, font="Calibri"):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color, font=font)
    return p


def rect(slide, l, t, w, h, fill):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s


def round_rect(slide, l, t, w, h, fill):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    try:
        s.adjustments[0] = 0.08
    except Exception:
        pass
    return s


def footer(slide, page):
    rect(slide, Inches(0.7), Inches(7.15), Inches(11.9), Pt(1.5), TEAL)
    box = slide.shapes.add_textbox(Inches(0.7), Inches(7.2), Inches(8), Inches(0.28))
    add_text(box, "Tairo ampio  ·  Confidentiel investisseurs", size=10, color=SLATE)
    num = slide.shapes.add_textbox(Inches(11.5), Inches(7.2), Inches(1.2), Inches(0.28))
    add_text(num, f"{page:02d} / {TOTAL}", size=10, color=SLATE, align=PP_ALIGN.RIGHT)


def slide_bg(slide, color=BG):
    rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(7.5), color)


def title_block(slide, eyebrow, title, subtitle=None):
    rect(slide, Inches(0), Inches(0), Inches(0.12), Inches(7.5), TEAL)
    eb = slide.shapes.add_textbox(Inches(0.7), Inches(0.35), Inches(11.5), Inches(0.35))
    add_text(eb, eyebrow.upper(), size=11, bold=True, color=TEAL)
    ti = slide.shapes.add_textbox(Inches(0.7), Inches(0.65), Inches(11.5), Inches(0.7))
    add_text(ti, title, size=32, bold=True, color=INK)
    if subtitle:
        su = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(11.5), Inches(0.4))
        add_text(su, subtitle, size=14, color=SLATE)


def card(slide, l, t, w, h, fill=WHITE):
    return round_rect(slide, l, t, w, h, fill)


# ═══════════════════════════════════════════════════════════════════════════════
# 01 — Title
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s, TEAL)
rect(s, Inches(0), Inches(6.4), Inches(13.333), Inches(1.1), TEAL_DARK)
mark = round_rect(s, Inches(0.9), Inches(1.5), Inches(0.7), Inches(0.7), WHITE)
mt = s.shapes.add_textbox(Inches(0.9), Inches(1.55), Inches(0.7), Inches(0.6))
add_text(mt, "T", size=28, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
brand = s.shapes.add_textbox(Inches(0.9), Inches(2.5), Inches(11), Inches(0.9))
add_text(brand, "Tairo ampio", size=54, bold=True, color=WHITE)
tag = s.shapes.add_textbox(Inches(0.9), Inches(3.4), Inches(11), Inches(0.5))
add_text(tag, "Marketplace de services à Madagascar", size=22, color=PEACH)
sub = s.shapes.add_textbox(Inches(0.9), Inches(4.2), Inches(11), Inches(0.8))
tf = add_text(sub, "Présentation produit  ·  État actuel  ·  Process", size=16, color=WHITE)
add_para(tf, "Deck investisseurs  ·  2026", size=14, color=PEACH, space_before=8)
foot = s.shapes.add_textbox(Inches(0.9), Inches(6.65), Inches(8), Inches(0.4))
add_text(foot, "TairoTairo  ·  Confidentiel", size=12, color=PEACH)
pg = s.shapes.add_textbox(Inches(11.5), Inches(6.65), Inches(1.2), Inches(0.4))
add_text(pg, "01 / 15", size=12, color=PEACH, align=PP_ALIGN.RIGHT)

# ═══════════════════════════════════════════════════════════════════════════════
# 02 — Sommaire
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
title_block(s, "Agenda", "Sommaire", "5 blocs pour comprendre Tairo ampio")
items = [
    ("01", "Problème & opportunité", "03 – 05"),
    ("02", "Produit & rôles", "06 – 07"),
    ("03", "Process & paiements", "08 – 10"),
    ("04", "Business & SWOT", "11 – 13"),
    ("05", "Next & call to action", "14 – 15"),
]
for i, (num, label, pages) in enumerate(items):
    y = Inches(2.0) + Inches(i * 0.9)
    card(s, Inches(0.7), y, Inches(11.9), Inches(0.75), WHITE)
    nbox = s.shapes.add_textbox(Inches(1.0), y + Inches(0.15), Inches(1.0), Inches(0.45))
    add_text(nbox, num, size=24, bold=True, color=TEAL)
    lbox = s.shapes.add_textbox(Inches(2.2), y + Inches(0.2), Inches(7), Inches(0.4))
    add_text(lbox, label, size=18, bold=True, color=INK)
    pbox = s.shapes.add_textbox(Inches(10.2), y + Inches(0.22), Inches(2), Inches(0.4))
    add_text(pbox, pages, size=14, color=SLATE, align=PP_ALIGN.RIGHT)
footer(s, 2)

# ═══════════════════════════════════════════════════════════════════════════════
# 03 — Problème
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
title_block(s, "Contexte", "Le problème", "Le marché local des services reste fragmenté et peu fiable")
problems = [
    ("Fragmentation", "Annonces Facebook, bouche-à-oreille, pas de place unique"),
    ("Confiance faible", "Pas de KYC, pas de preuve de qualité, litiges fréquents"),
    ("Paiement hors app", "Cash / Mobile Money direct → zéro protection"),
    ("Pas de traçabilité", "Pas d'avis liés à une prestation réellement payée"),
]
for i, (t, d) in enumerate(problems):
    col = i % 2
    row = i // 2
    x = Inches(0.7) + Inches(col * 6.15)
    y = Inches(2.1) + Inches(row * 2.2)
    card(s, x, y, Inches(5.9), Inches(1.95), WHITE)
    rect(s, x, y, Inches(0.12), Inches(1.95), TEAL if i % 2 == 0 else PEACH)
    tb = s.shapes.add_textbox(x + Inches(0.4), y + Inches(0.4), Inches(5.2), Inches(0.45))
    add_text(tb, t, size=20, bold=True, color=INK)
    db = s.shapes.add_textbox(x + Inches(0.4), y + Inches(1.0), Inches(5.2), Inches(0.7))
    add_text(db, d, size=14, color=SLATE)
footer(s, 3)

# ═══════════════════════════════════════════════════════════════════════════════
# 04 — Solution
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
title_block(s, "Proposition", "La solution Tairo ampio", "Une marketplace sécurisée, de bout en bout")
pillars = [
    ("01", "Marketplace 2 faces", "Annonces prestataires + demandes clients"),
    ("02", "Paiement séquestre", "Fonds bloqués jusqu'à validation client"),
    ("03", "KYC prestataires", "Identité vérifiée avant d'opérer"),
    ("04", "Avis crédibles", "Notes liées aux paiements in-app uniquement"),
]
for i, (n, t, d) in enumerate(pillars):
    x = Inches(0.7) + Inches(i * 3.1)
    card(s, x, Inches(2.2), Inches(2.95), Inches(4.3), WHITE)
    nb = s.shapes.add_textbox(x + Inches(0.25), Inches(2.5), Inches(2.4), Inches(0.5))
    add_text(nb, n, size=28, bold=True, color=TEAL)
    tb = s.shapes.add_textbox(x + Inches(0.25), Inches(3.3), Inches(2.4), Inches(1.0))
    add_text(tb, t, size=16, bold=True, color=INK)
    db = s.shapes.add_textbox(x + Inches(0.25), Inches(4.5), Inches(2.4), Inches(1.5))
    add_text(db, d, size=13, color=SLATE)
footer(s, 4)

# ═══════════════════════════════════════════════════════════════════════════════
# 05 — Pourquoi maintenant
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
title_block(s, "Timing", "Pourquoi maintenant", "Trois vents favorables à Madagascar")
winds = [
    ("Mobile Money mature", "Orange Money · MVola · Airtel Money", "Rails de paiement déjà massifs"),
    ("Demande urbaine", "Services du quotidien en forte", "demande (Tana & grandes villes)"),
    ("Gap trust & escrow", "Aucun acteur local n'offre encore", "un séquestre + KYC + avis liés"),
]
for i, (t, d1, d2) in enumerate(winds):
    x = Inches(0.7) + Inches(i * 4.15)
    card(s, x, Inches(2.3), Inches(3.95), Inches(3.8), WHITE)
    rect(s, x, Inches(2.3), Inches(3.95), Inches(0.12), TEAL)
    tb = s.shapes.add_textbox(x + Inches(0.3), Inches(2.8), Inches(3.35), Inches(0.8))
    add_text(tb, t, size=18, bold=True, color=INK)
    db = s.shapes.add_textbox(x + Inches(0.3), Inches(3.8), Inches(3.35), Inches(1.8))
    tf = add_text(db, d1, size=14, color=SLATE)
    add_para(tf, d2, size=14, color=SLATE, space_before=6)
footer(s, 5)

# ═══════════════════════════════════════════════════════════════════════════════
# 06 — Produit + donut
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
title_block(s, "Produit", "Produit en un regard", "3 rôles · 14 catégories · web responsive")
roles = [
    ("CLIENT", "Réserver, demander, payer, noter"),
    ("PRESTATAIRE", "Annoncer, proposer, exécuter, encaisser"),
    ("ADMIN", "KYC, modération, stats, spotlight"),
]
for i, (t, d) in enumerate(roles):
    y = Inches(2.1) + Inches(i * 1.45)
    card(s, Inches(0.7), y, Inches(6.2), Inches(1.3), WHITE)
    rect(s, Inches(0.7), y, Inches(0.12), Inches(1.3), TEAL)
    tb = s.shapes.add_textbox(Inches(1.1), y + Inches(0.25), Inches(5.5), Inches(0.4))
    add_text(tb, t, size=16, bold=True, color=TEAL)
    db = s.shapes.add_textbox(Inches(1.1), y + Inches(0.7), Inches(5.5), Inches(0.4))
    add_text(db, d, size=13, color=SLATE)

chart_data = CategoryChartData()
chart_data.categories = ["Clients", "Prestataires", "Admin"]
chart_data.add_series("Rôles", (55, 40, 5))
chart = s.shapes.add_chart(
    XL_CHART_TYPE.DOUGHNUT, Inches(7.3), Inches(2.1), Inches(5.3), Inches(4.4), chart_data
).chart
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.include_in_layout = False
cap = s.shapes.add_textbox(Inches(7.3), Inches(6.55), Inches(5.3), Inches(0.35))
add_text(cap, "Répartition cible des rôles (illustratif)  ·  Source : modèle produit", size=9, color=SLATE, align=PP_ALIGN.CENTER)
footer(s, 6)

# ═══════════════════════════════════════════════════════════════════════════════
# 07 — Ce qui est live
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
title_block(s, "État produit", "Ce qui est live aujourd'hui", "MVP feature-complete — prêt pour l'intégration Mobile Money")
features = [
    "Auth & sécurité (JWT, OTP, CSRF)",
    "Annonces & demandes",
    "Réservations & planning",
    "Messagerie temps réel + négo prix",
    "Escrow (séquestre) in-app",
    "KYC prestataires (admin)",
    "Avis liés aux paiements",
    "Notifications (in-app, email, push)",
    "Portfolio prestataire",
    "Abonnements & spotlight",
    "Facture PDF auto",
    "Dashboard admin & exports",
]
for i, f in enumerate(features):
    col = i % 3
    row = i // 3
    x = Inches(0.7) + Inches(col * 4.15)
    y = Inches(2.1) + Inches(row * 1.15)
    card(s, x, y, Inches(3.95), Inches(1.0), WHITE)
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.2), y + Inches(0.3), Inches(0.35), Inches(0.35))
    c.fill.solid()
    c.fill.fore_color.rgb = TEAL_SOFT
    c.line.fill.background()
    ck = s.shapes.add_textbox(x + Inches(0.2), y + Inches(0.3), Inches(0.35), Inches(0.35))
    add_text(ck, "✓", size=12, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    tb = s.shapes.add_textbox(x + Inches(0.7), y + Inches(0.3), Inches(3.0), Inches(0.45))
    add_text(tb, f, size=12, bold=True, color=INK)
footer(s, 7)

# ═══════════════════════════════════════════════════════════════════════════════
# 08 — Parcours client
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
title_block(s, "Process", "Parcours client", "De la découverte à l'avis — en 6 étapes")
steps = [
    ("1", "Découvrir", "Services ou", "demande"),
    ("2", "Réserver", "Service ou", "proposition"),
    ("3", "Payer", "Séquestre", "Mobile Money"),
    ("4", "Recevoir", "Prestation", "en cours"),
    ("5", "Valider", "Libération", "des fonds"),
    ("6", "Noter", "Avis + facture", "PDF"),
]
for i, (n, t, d1, d2) in enumerate(steps):
    x = Inches(0.55) + Inches(i * 2.15)
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.55), Inches(2.5), Inches(0.7), Inches(0.7))
    c.fill.solid()
    c.fill.fore_color.rgb = TEAL
    c.line.fill.background()
    nb = s.shapes.add_textbox(x + Inches(0.55), Inches(2.6), Inches(0.7), Inches(0.55))
    add_text(nb, n, size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < 5:
        rect(s, x + Inches(1.4), Inches(2.8), Inches(1.1), Pt(3), TEAL_SOFT)
    tb = s.shapes.add_textbox(x, Inches(3.5), Inches(1.9), Inches(0.4))
    add_text(tb, t, size=14, bold=True, color=INK, align=PP_ALIGN.CENTER)
    db = s.shapes.add_textbox(x, Inches(4.0), Inches(1.9), Inches(1.0))
    tf = add_text(db, d1, size=11, color=SLATE, align=PP_ALIGN.CENTER)
    add_para(tf, d2, size=11, color=SLATE, align=PP_ALIGN.CENTER, space_before=2)

card(s, Inches(0.7), Inches(5.4), Inches(11.9), Inches(1.2), TEAL_SOFT)
hb = s.shapes.add_textbox(Inches(1.0), Inches(5.65), Inches(11.3), Inches(0.7))
tf = add_text(hb, "Deux chemins d'entrée", size=14, bold=True, color=TEAL)
add_para(tf, "A · Réserver une annonce   ·   B · Publier une demande et accepter une proposition", size=13, color=INK, space_before=4)
footer(s, 8)

# ═══════════════════════════════════════════════════════════════════════════════
# 09 — Parcours prestataire
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
title_block(s, "Process", "Parcours prestataire", "KYC d'abord — puis activité et encaissement")
steps = [
    ("1", "S'inscrire", "Compte", "prestataire"),
    ("2", "KYC", "CIN →", "validation admin"),
    ("3", "Publier", "Annonces &", "propositions"),
    ("4", "Exécuter", "Démarrer →", "marquer terminé"),
    ("5", "Encaisser", "Versement après", "validation client"),
    ("6", "Booster", "Abo spotlight", "& portfolio"),
]
for i, (n, t, d1, d2) in enumerate(steps):
    x = Inches(0.55) + Inches(i * 2.15)
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.55), Inches(2.5), Inches(0.7), Inches(0.7))
    c.fill.solid()
    c.fill.fore_color.rgb = TEAL if i != 1 else AMBER
    c.line.fill.background()
    nb = s.shapes.add_textbox(x + Inches(0.55), Inches(2.6), Inches(0.7), Inches(0.55))
    add_text(nb, n, size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < 5:
        rect(s, x + Inches(1.4), Inches(2.8), Inches(1.1), Pt(3), TEAL_SOFT)
    tb = s.shapes.add_textbox(x, Inches(3.5), Inches(1.9), Inches(0.4))
    add_text(tb, t, size=14, bold=True, color=INK, align=PP_ALIGN.CENTER)
    db = s.shapes.add_textbox(x, Inches(4.0), Inches(1.9), Inches(1.0))
    tf = add_text(db, d1, size=11, color=SLATE, align=PP_ALIGN.CENTER)
    add_para(tf, d2, size=11, color=SLATE, align=PP_ALIGN.CENTER, space_before=2)

card(s, Inches(0.7), Inches(5.4), Inches(11.9), Inches(1.2), PEACH_SOFT)
hb = s.shapes.add_textbox(Inches(1.0), Inches(5.65), Inches(11.3), Inches(0.7))
tf = add_text(hb, "Gate KYC", size=14, bold=True, color=AMBER)
add_para(tf, "Sans KYC approuvé : pas d'annonces, pas de réponses aux demandes, pas de spotlight", size=13, color=INK, space_before=4)
footer(s, 9)

# ═══════════════════════════════════════════════════════════════════════════════
# 10 — Process paiement
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
title_block(s, "Paiements", "Process paiement (escrow)", "Fonds protégés jusqu'à validation — Mobile Money réel = next")

chart_data = CategoryChartData()
chart_data.categories = ["Confirmé", "Payé (séquestre)", "En cours", "Terminé", "Validé / libéré"]
chart_data.add_series("Tunnel", (100, 85, 80, 75, 70))
chart = s.shapes.add_chart(
    XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.5), Inches(2.0), Inches(7.2), Inches(4.5), chart_data
).chart
chart.has_legend = False

statuses = [
    ("ESCROWED", "Fonds bloqués Tairo ampio", TEAL),
    ("RELEASED", "Versement prestataire", GREEN),
    ("REFUNDED", "Remboursement client", AMBER),
]
for i, (st, desc, col) in enumerate(statuses):
    y = Inches(2.2) + Inches(i * 1.4)
    card(s, Inches(8.0), y, Inches(4.6), Inches(1.2), WHITE)
    rect(s, Inches(8.0), y, Inches(0.12), Inches(1.2), col)
    tb = s.shapes.add_textbox(Inches(8.4), y + Inches(0.25), Inches(4.0), Inches(0.35))
    add_text(tb, st, size=14, bold=True, color=col)
    db = s.shapes.add_textbox(Inches(8.4), y + Inches(0.65), Inches(4.0), Inches(0.35))
    add_text(db, desc, size=12, color=SLATE)

cap = s.shapes.add_textbox(Inches(0.7), Inches(6.55), Inches(7), Inches(0.3))
add_text(cap, "Tunnel escrow (illustratif)  ·  Source : workflow produit", size=9, color=SLATE)
footer(s, 10)

# ═══════════════════════════════════════════════════════════════════════════════
# 11 — Monétisation
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
title_block(s, "Business", "Monétisation", "Revenus récurrents + take-rate escrow (roadmap)")

card(s, Inches(0.7), Inches(2.1), Inches(5.5), Inches(4.4), WHITE)
nb = s.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(5.0), Inches(0.4))
add_text(nb, "ABONNEMENT PRESTATAIRE", size=12, bold=True, color=TEAL)
big = s.shapes.add_textbox(Inches(1.0), Inches(3.0), Inches(5.0), Inches(1.0))
add_text(big, "75 000 Ar", size=44, bold=True, color=INK)
sub = s.shapes.add_textbox(Inches(1.0), Inches(4.1), Inches(5.0), Inches(0.4))
add_text(sub, "/ mois  ·  plan 3 mois −10 %", size=14, color=SLATE)
benefits = ["· Profil mis en avant (accueil)", "· Suggestions recherche", "· 1 annonce featured"]
bb = s.shapes.add_textbox(Inches(1.0), Inches(4.7), Inches(5.0), Inches(1.5))
tf = add_text(bb, benefits[0], size=13, color=INK)
for b in benefits[1:]:
    add_para(tf, b, size=13, color=INK, space_before=4)

chart_data = CategoryChartData()
chart_data.categories = ["1 mois", "3 mois"]
chart_data.add_series("Prix (Ar)", (75000, 202500))
chart = s.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(6.6), Inches(2.1), Inches(6.0), Inches(4.0), chart_data
).chart
chart.has_legend = False
cap = s.shapes.add_textbox(Inches(6.6), Inches(6.2), Inches(6.0), Inches(0.35))
add_text(cap, "Comparaison plans abo  ·  Source : pricing produit", size=9, color=SLATE, align=PP_ALIGN.CENTER)
note = s.shapes.add_textbox(Inches(0.7), Inches(6.55), Inches(11.9), Inches(0.3))
add_text(note, "Roadmap : commission escrow sur transactions  ·  Take-rate à définir avec l'intégration MM", size=11, color=SLATE)
footer(s, 11)

# ═══════════════════════════════════════════════════════════════════════════════
# 12 — Stack
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
title_block(s, "Tech", "Stack & opérations", "Architecture prête pour un VPS long-lived")
stack = [
    ("Frontend", "Next.js 16 · React 19 · Tailwind"),
    ("Backend", "Route Handlers · Node HTTP custom"),
    ("Data", "PostgreSQL · Prisma · RLS"),
    ("Realtime", "WebSockets · Redis pub/sub"),
    ("Media", "Sharp (WebP) · stockage local"),
    ("Ops", "Docker · healthcheck · cron abo"),
]
for i, (t, d) in enumerate(stack):
    col = i % 3
    row = i // 3
    x = Inches(0.7) + Inches(col * 4.15)
    y = Inches(2.2) + Inches(row * 2.15)
    card(s, x, y, Inches(3.95), Inches(1.9), WHITE)
    tb = s.shapes.add_textbox(x + Inches(0.3), y + Inches(0.35), Inches(3.35), Inches(0.4))
    add_text(tb, t, size=14, bold=True, color=TEAL)
    db = s.shapes.add_textbox(x + Inches(0.3), y + Inches(0.9), Inches(3.35), Inches(0.7))
    add_text(db, d, size=13, color=INK)
footer(s, 12)

# ═══════════════════════════════════════════════════════════════════════════════
# 13 — SWOT
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
title_block(s, "Analyse", "Matrice SWOT", "Forces, faiblesses, opportunités, menaces")
swot = [
    ("Forces", GREEN, GREEN_SOFT, [
        "MVP feature-complete",
        "Escrow + KYC + avis liés",
        "Messagerie + négo prix",
        "Admin ops & exports",
        "Stack moderne & testée",
    ]),
    ("Faiblesses", AMBER, AMBER_SOFT, [
        "Mobile Money non branché",
        "Payouts encore manuels",
        "Stockage fichiers local",
        "Pas d'app native",
        "Traction à prouver",
    ]),
    ("Opportunités", TEAL, TEAL_SOFT, [
        "Rails MM déjà massifs",
        "Demande services urbains",
        "Take-rate escrow",
        "B2B / entreprises",
        "Expansion villes MG",
    ]),
    ("Menaces", RED, RED_SOFT, [
        "Groupes Facebook gratuits",
        "Concurrents MM-first",
        "Régulation paiements",
        "Adoption prestataires",
        "Confiance marché",
    ]),
]
for i, (title, accent, soft, bullets) in enumerate(swot):
    col = i % 2
    row = i // 2
    x = Inches(0.7) + Inches(col * 6.2)
    y = Inches(1.95) + Inches(row * 2.45)
    card(s, x, y, Inches(5.95), Inches(2.3), soft)
    rect(s, x, y, Inches(5.95), Inches(0.08), accent)
    tb = s.shapes.add_textbox(x + Inches(0.3), y + Inches(0.2), Inches(5.3), Inches(0.35))
    add_text(tb, title.upper(), size=13, bold=True, color=accent)
    bb = s.shapes.add_textbox(x + Inches(0.3), y + Inches(0.6), Inches(5.3), Inches(1.55))
    tf = add_text(bb, "·  " + bullets[0], size=12, color=INK)
    for b in bullets[1:]:
        add_para(tf, "·  " + b, size=12, color=INK, space_before=3)
footer(s, 13)

# ═══════════════════════════════════════════════════════════════════════════════
# 14 — État & next
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
title_block(s, "Roadmap", "État actuel & next", "Priorité claire : brancher le rail de paiement réel")

card(s, Inches(0.7), Inches(2.1), Inches(5.9), Inches(4.4), WHITE)
rect(s, Inches(0.7), Inches(2.1), Inches(5.9), Inches(0.55), TEAL)
ht = s.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(5.3), Inches(0.4))
add_text(ht, "AUJOURD'HUI", size=14, bold=True, color=WHITE)
now = [
    "Marketplace 2 faces opérationnelle",
    "Workflow escrow complet (logique)",
    "KYC + admin + factures PDF",
    "Messagerie temps réel",
    "Abonnements & spotlight",
]
nb = s.shapes.add_textbox(Inches(1.0), Inches(2.9), Inches(5.3), Inches(3.2))
tf = add_text(nb, "✓  " + now[0], size=14, color=INK)
for item in now[1:]:
    add_para(tf, "✓  " + item, size=14, color=INK, space_before=10)

card(s, Inches(6.85), Inches(2.1), Inches(5.9), Inches(4.4), WHITE)
rect(s, Inches(6.85), Inches(2.1), Inches(5.9), Inches(0.55), PEACH)
ht = s.shapes.add_textbox(Inches(7.15), Inches(2.2), Inches(5.3), Inches(0.4))
add_text(ht, "NEXT", size=14, bold=True, color=INK)
nxt = [
    "#1  API Mobile Money (capture + payout)",
    "Versements auto prestataires",
    "Object storage (S3 / R2)",
    "Scale multi-instances + Redis",
    "Acquisition & métriques live",
]
nb = s.shapes.add_textbox(Inches(7.15), Inches(2.9), Inches(5.3), Inches(3.2))
tf = add_text(nb, "→  " + nxt[0], size=14, bold=True, color=TEAL)
for item in nxt[1:]:
    add_para(tf, "→  " + item, size=14, color=INK, space_before=10)
footer(s, 14)

# ═══════════════════════════════════════════════════════════════════════════════
# 15 — CTA
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s, TEAL)
rect(s, Inches(0), Inches(6.4), Inches(13.333), Inches(1.1), TEAL_DARK)

t = s.shapes.add_textbox(Inches(0.9), Inches(1.8), Inches(11.5), Inches(1.4))
tf = add_text(t, "Passons à l'intégration", size=36, bold=True, color=WHITE)
add_para(tf, "Mobile Money.", size=36, bold=True, color=PEACH, space_before=4)

sub = s.shapes.add_textbox(Inches(0.9), Inches(3.6), Inches(11), Inches(1.0))
tf = add_text(sub, "MVP prêt. Rail de paiement = levier de scale.", size=18, color=WHITE)
add_para(tf, "Discutons investissement, partenariats MM et go-to-market.", size=16, color=PEACH, space_before=10)

pill = round_rect(s, Inches(0.9), Inches(5.1), Inches(4.2), Inches(0.7), WHITE)
cta = s.shapes.add_textbox(Inches(0.9), Inches(5.22), Inches(4.2), Inches(0.5))
add_text(cta, "Parlons-en →", size=18, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

foot = s.shapes.add_textbox(Inches(0.9), Inches(6.65), Inches(8), Inches(0.4))
add_text(foot, "Tairo ampio  ·  TairoTairo  ·  Madagascar", size=12, color=PEACH)
pg = s.shapes.add_textbox(Inches(11.5), Inches(6.65), Inches(1.2), Inches(0.4))
add_text(pg, "15 / 15", size=12, color=PEACH, align=PP_ALIGN.RIGHT)

out = "/Users/macbookair/Desktop/tairotairo/presentations/Tairo_ampio_Investor_Deck.pptx"
os.makedirs(os.path.dirname(out), exist_ok=True)
prs.save(out)
print("OK", out, os.path.getsize(out))
