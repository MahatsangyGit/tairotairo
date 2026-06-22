#!/usr/bin/env python3
"""Génère la présentation PowerPoint du projet Tairo ampio."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# Charte graphique — app/lib/site.ts
BRAND_PRIMARY = "#069494"
BRAND_SECONDARY = "#ffffff"
BRAND_TERTIARY = "#FAD0C9"
INK = "#1a2e2e"
MUTED = "#4a6363"
LIGHT_BG = "#f0fafa"

OUTPUT = Path(__file__).resolve().parent.parent / "docs" / "Tairo-ampio-presentation.pptx"


def hex_rgb(h: str) -> RGBColor:
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


PRIMARY = hex_rgb(BRAND_PRIMARY)
SECONDARY = hex_rgb(BRAND_SECONDARY)
TERTIARY = hex_rgb(BRAND_TERTIARY)
INK_COLOR = hex_rgb(INK)
MUTED_COLOR = hex_rgb(MUTED)
LIGHT = hex_rgb(LIGHT_BG)


def set_slide_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_rgb, line_rgb=None):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, size=18, bold=False, color=INK_COLOR, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return box


def add_bullets(slide, left, top, width, height, title, items, title_size=28, bullet_size=16):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(title_size)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p.space_after = Pt(12)
    for item in items:
        bp = tf.add_paragraph()
        bp.text = item
        bp.level = 0
        bp.font.size = Pt(bullet_size)
        bp.font.color.rgb = INK_COLOR
        bp.space_after = Pt(6)
    return box


def add_header_bar(slide, title: str, subtitle: str = ""):
    add_rect(slide, Inches(0), Inches(0), Inches(13.33), Inches(1.1), PRIMARY)
    add_text_box(
        slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.55),
        title, size=32, bold=True, color=SECONDARY,
    )
    if subtitle:
        add_text_box(
            slide, Inches(0.5), Inches(0.62), Inches(12), Inches(0.4),
            subtitle, size=14, color=SECONDARY,
        )


def add_footer(slide, text: str = "Tairo ampio · TairoTairo · Marketplace Madagascar"):
    add_rect(slide, Inches(0), Inches(7.0), Inches(13.33), Inches(0.5), PRIMARY)
    add_text_box(
        slide, Inches(0.5), Inches(7.05), Inches(12), Inches(0.35),
        text, size=10, color=SECONDARY, align=PP_ALIGN.CENTER,
    )


def add_uml_box(slide, left, top, width, height, title, lines, fill=TERTIARY):
    shape = add_rect(slide, left, top, width, height, fill, PRIMARY)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = INK_COLOR
    p.alignment = PP_ALIGN.CENTER
    for line in lines:
        lp = tf.add_paragraph()
        lp.text = line
        lp.font.size = Pt(9)
        lp.font.color.rgb = MUTED_COLOR
        lp.alignment = PP_ALIGN.CENTER
    return shape


def connect_boxes(slide, x1, y1, x2, y2):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = PRIMARY
    conn.line.width = Pt(1.5)
    return conn


def slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PRIMARY)
    add_rect(slide, Inches(0), Inches(5.8), Inches(13.33), Inches(1.7), TERTIARY)
    add_text_box(
        slide, Inches(0.8), Inches(1.8), Inches(11), Inches(1.2),
        "Tairo ampio", size=54, bold=True, color=SECONDARY,
    )
    add_text_box(
        slide, Inches(0.8), Inches(3.0), Inches(11), Inches(0.8),
        "Marketplace de services à Madagascar", size=28, color=SECONDARY,
    )
    add_text_box(
        slide, Inches(0.8), Inches(4.0), Inches(11), Inches(0.5),
        "Présentation technique du projet · État actuel & feuille de route", size=16, color=SECONDARY,
    )
    add_text_box(
        slide, Inches(0.8), Inches(6.1), Inches(11), Inches(0.5),
        "TairoTairo · Mai 2026", size=14, color=INK_COLOR,
    )


def slide_vision(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, SECONDARY)
    add_header_bar(slide, "Vision du projet")
    add_bullets(
        slide, Inches(0.6), Inches(1.4), Inches(5.8), Inches(5.5),
        "Objectif",
        [
            "Mettre en relation clients et prestataires de services à Madagascar",
            "Offrir annonces (services), demandes (requests), réservations et messagerie",
            "Vérification KYC des prestataires et abonnements pour la visibilité",
            "Paiements Mobile Money (Orange, MVola, Airtel) — en cours d'intégration",
        ],
    )
    add_bullets(
        slide, Inches(6.8), Inches(1.4), Inches(6), Inches(5.5),
        "Acteurs",
        [
            "Client — publie des demandes, réserve, note, discute",
            "Prestataire — publie des services, portfolio, KYC, abonnement",
            "Admin — modération KYC, spotlight, stats, export",
        ],
        bullet_size=15,
    )
    add_footer(slide)


def slide_stack(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT)
    add_header_bar(slide, "Stack technique")
    stacks = [
        ("Frontend", ["Next.js 16.2 · App Router", "React 19 · TypeScript", "Tailwind CSS v4", "shadcn/ui · Lucide · next-themes (dark)"]),
        ("Backend", ["Route Handlers API (60 routes)", "Prisma 7 + PostgreSQL", "JWT custom (jsonwebtoken + jose)", "bcryptjs · nodemailer · web-push"]),
        ("Stockage", ["Fichiers locaux storage/", "Avatars, KYC, portfolio, couvertures annonces"]),
        ("Qualité", ["ESLint · Playwright E2E", "Migrations SQL Supabase"]),
    ]
    x_positions = [Inches(0.4), Inches(3.5), Inches(6.6), Inches(9.7)]
    for i, (title, items) in enumerate(stacks):
        add_rect(slide, x_positions[i], Inches(1.5), Inches(2.9), Inches(5.2), SECONDARY, PRIMARY)
        add_text_box(slide, x_positions[i] + Inches(0.15), Inches(1.65), Inches(2.6), Inches(0.4), title, size=18, bold=True, color=PRIMARY)
        y = Inches(2.2)
        for item in items:
            add_text_box(slide, x_positions[i] + Inches(0.15), y, Inches(2.6), Inches(0.9), f"• {item}", size=12, color=INK_COLOR)
            y += Inches(0.85)
    add_footer(slide)


def slide_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, SECONDARY)
    add_header_bar(slide, "Architecture applicative")
    layers = [
        ("Navigateur", "Pages publiques · Dashboard client/prestataire · Admin"),
        ("Next.js App Router", "RSC + Client Components · Middleware auth"),
        ("API Routes /api/*", "60 endpoints REST · Validation · Prisma"),
        ("PostgreSQL", "18 modèles Prisma · Migrations Supabase"),
        ("Services externes", "SMTP · Web Push (VAPID) · Cron abonnements"),
    ]
    y = Inches(1.5)
    for i, (name, desc) in enumerate(layers):
        fill = TERTIARY if i % 2 == 0 else LIGHT
        add_rect(slide, Inches(1.5), y, Inches(10.3), Inches(0.95), fill, PRIMARY)
        add_text_box(slide, Inches(1.7), y + Inches(0.1), Inches(3), Inches(0.4), name, size=16, bold=True, color=PRIMARY)
        add_text_box(slide, Inches(4.8), y + Inches(0.15), Inches(6.5), Inches(0.5), desc, size=13, color=INK_COLOR)
        if i < len(layers) - 1:
            add_text_box(slide, Inches(6.2), y + Inches(0.85), Inches(0.5), Inches(0.3), "▼", size=12, color=PRIMARY, align=PP_ALIGN.CENTER)
        y += Inches(1.1)
    add_footer(slide)


def slide_prisma_uml(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, SECONDARY)
    add_header_bar(slide, "Modèle de données (UML simplifié)", "18 entités Prisma · PostgreSQL")
    # User central
    add_uml_box(slide, Inches(5.4), Inches(1.5), Inches(2.5), Inches(0.9), "User", ["CLIENT · PROVIDER · ADMIN"])
    entities = [
        (Inches(0.5), Inches(2.8), "Service", ["providerId"]),
        (Inches(3.0), Inches(2.8), "ServiceRequest", ["clientId"]),
        (Inches(5.5), Inches(2.8), "Booking", ["client · provider"]),
        (Inches(8.0), Inches(2.8), "Conversation", ["client · provider"]),
        (Inches(10.5), Inches(2.8), "Message", ["TEXT · PRICE_OFFER"]),
        (Inches(0.5), Inches(4.2), "Review", ["bookingId"]),
        (Inches(3.0), Inches(4.2), "RequestResponse", ["PENDING→COMPLETED"]),
        (Inches(5.5), Inches(4.2), "Transaction", ["MGA · Mobile Money"]),
        (Inches(8.0), Inches(4.2), "ProviderKycDocument", ["CIN recto/verso"]),
        (Inches(10.5), Inches(4.2), "ProviderPortfolioItem", ["images + comments"]),
        (Inches(2.0), Inches(5.6), "ProviderSubscription", ["expiresAt"]),
        (Inches(5.0), Inches(5.6), "Notification", ["email + push"]),
        (Inches(8.0), Inches(5.6), "PushSubscription", ["web-push"]),
    ]
    for left, top, title, lines in entities:
        add_uml_box(slide, left, top, Inches(2.2), Inches(0.85), title, lines, LIGHT if title in ("Transaction",) else TERTIARY)
    connect_boxes(slide, Inches(6.65), Inches(2.4), Inches(1.6), Inches(2.8))
    connect_boxes(slide, Inches(6.65), Inches(2.4), Inches(6.6), Inches(2.8))
    connect_boxes(slide, Inches(6.65), Inches(2.4), Inches(11.6), Inches(2.8))
    add_text_box(
        slide, Inches(0.5), Inches(6.55), Inches(12), Inches(0.35),
        "Transaction : modèle présent mais non alimenté par les flux de paiement réels",
        size=11, color=MUTED_COLOR,
    )
    add_footer(slide)


def slide_api_overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT)
    add_header_bar(slide, "Routes API — Vue d'ensemble", "60 endpoints REST sous /api")
    domains = [
        ("Auth (8)", "login · register · logout · me · OTP email · reset password"),
        ("Users (3)", "me · avatar · [id]/avatar"),
        ("Services (3)", "CRUD · cover image"),
        ("Requests (5)", "CRUD · responses · cover"),
        ("Bookings (2)", "liste · [id] statut"),
        ("Reviews (1)", "création après réservation"),
        ("Providers (2)", "profil public · portfolio"),
        ("Provider (11)", "KYC · portfolio · subscription · featured"),
        ("Conversations (8)", "messagerie · offres de prix · unread"),
        ("Notifications (6)", "liste · préférences · push subscribe"),
        ("Admin (10)", "users · KYC · stats · spotlight · export"),
        ("Cron (1)", "expire-subscriptions"),
    ]
    col1, col2 = Inches(0.5), Inches(6.9)
    for i, (name, routes) in enumerate(domains):
        col = col1 if i < 6 else col2
        row = i if i < 6 else i - 6
        top = Inches(1.35) + row * Inches(0.92)
        add_rect(slide, col, top, Inches(6.2), Inches(0.82), SECONDARY, PRIMARY)
        add_text_box(slide, col + Inches(0.12), top + Inches(0.08), Inches(1.6), Inches(0.35), name, size=13, bold=True, color=PRIMARY)
        add_text_box(slide, col + Inches(1.7), top + Inches(0.12), Inches(4.3), Inches(0.55), routes, size=10, color=INK_COLOR)
    add_footer(slide)


def slide_api_uml_auth(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, SECONDARY)
    add_header_bar(slide, "UML — Domaine Auth & Utilisateurs")
    client = add_uml_box(slide, Inches(0.6), Inches(2.0), Inches(2.0), Inches(0.7), "Client", ["Browser / App"])
    api = add_uml_box(slide, Inches(4.0), Inches(1.8), Inches(5.3), Inches(4.5), "/api/auth · /api/users", [
        "POST /auth/login",
        "POST /auth/register",
        "POST /auth/logout",
        "GET  /auth/me",
        "POST /auth/email/send-otp",
        "POST /auth/email/verify-otp",
        "POST /auth/forgot-password",
        "POST /auth/reset-password",
        "GET/PATCH /users/me",
        "POST/GET/DELETE /users/me/avatar",
        "GET /users/[id]/avatar",
    ], LIGHT)
    db = add_uml_box(slide, Inches(10.2), Inches(2.0), Inches(2.5), Inches(1.2), "PostgreSQL", ["User", "EmailOtp", "PasswordResetToken"])
    connect_boxes(slide, Inches(2.6), Inches(2.35), Inches(4.0), Inches(2.35))
    connect_boxes(slide, Inches(9.3), Inches(2.8), Inches(10.2), Inches(2.6))
    add_text_box(slide, Inches(0.6), Inches(5.8), Inches(12), Inches(0.8),
                 "JWT en cookie httpOnly · Rôles CLIENT / PROVIDER / ADMIN · Vérification email OTP",
                 size=13, color=MUTED_COLOR)
    add_footer(slide)


def slide_api_uml_marketplace(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, SECONDARY)
    add_header_bar(slide, "UML — Marketplace (Services & Demandes)")
    add_uml_box(slide, Inches(0.5), Inches(1.5), Inches(3.8), Inches(2.2), "/api/services", [
        "GET/POST /services",
        "GET/PATCH/DELETE /services/[id]",
        "POST/GET/DELETE /services/[id]/cover",
    ])
    add_uml_box(slide, Inches(4.7), Inches(1.5), Inches(4.0), Inches(2.8), "/api/requests", [
        "GET/POST /requests",
        "GET/PATCH/DELETE /requests/[id]",
        "POST/GET /requests/[id]/responses",
        "PATCH /requests/[id]/responses/[responseId]",
        "POST/GET/DELETE /requests/[id]/cover",
        "POST /responses (legacy)",
    ])
    add_uml_box(slide, Inches(9.0), Inches(1.5), Inches(3.8), Inches(1.6), "/api/bookings · /api/reviews", [
        "GET/POST /bookings",
        "PATCH /bookings/[id]",
        "POST /reviews",
    ])
    connect_boxes(slide, Inches(2.4), Inches(3.7), Inches(4.7), Inches(3.2))
    connect_boxes(slide, Inches(8.7), Inches(3.0), Inches(9.0), Inches(2.5))
    add_uml_box(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(1.8), "Flux métier", [
        "Service → Booking direct · Request → Response → Booking",
        "Couverture image (1 max) · SEO catégories /services/categorie/[slug]",
        "Affichage public · Dashboard client & prestataire",
    ], LIGHT)
    add_footer(slide)


def slide_api_uml_messaging(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT)
    add_header_bar(slide, "UML — Messagerie & Négociation")
    add_uml_box(slide, Inches(0.5), Inches(1.5), Inches(3.5), Inches(3.5), "/api/conversations", [
        "GET/POST /conversations",
        "POST /conversations/open",
        "GET /conversations/unread-count",
        "GET/PATCH /conversations/[id]",
        "GET/POST /conversations/[id]/messages",
        "GET/POST /conversations/[id]/price-offers",
        "POST .../price-offers/[messageId]/accept",
    ])
    add_uml_box(slide, Inches(4.5), Inches(1.5), Inches(3.8), Inches(2.0), "Message kinds", [
        "TEXT — message classique",
        "PRICE_OFFER — offre liée service ou response",
        "Statuts: PENDING · ACCEPTED · SUPERSEDED",
    ], TERTIARY)
    add_uml_box(slide, Inches(8.8), Inches(1.5), Inches(4.0), Inches(2.5), "/api/notifications", [
        "GET /notifications",
        "PATCH /notifications/[id]",
        "POST /notifications/read-all",
        "POST /notifications/clear-all",
        "GET/PATCH /notifications/preferences",
        "POST /notifications/push/subscribe",
    ])
    add_text_box(slide, Inches(0.5), Inches(5.5), Inches(12), Inches(1.2),
                 "Messagerie temps réel : polling actuel (pas de WebSocket)\n"
                 "Notifications email (nodemailer) + push navigateur (web-push / VAPID)",
                 size=14, color=INK_COLOR)
    add_footer(slide)


def slide_api_uml_provider(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, SECONDARY)
    add_header_bar(slide, "UML — Prestataire, KYC & Abonnement")
    boxes = [
        (Inches(0.4), "/api/provider/kyc", ["GET statut", "POST submit", "POST upload", "DELETE documents/[id]"]),
        (Inches(3.5), "/api/provider/portfolio", ["CRUD portfolio", "POST image", "GET/POST comments"]),
        (Inches(6.6), "/api/provider/subscription", ["GET statut", "POST purchase (simulé)", "Cron expire"]),
        (Inches(9.7), "/api/providers/[id]", ["Profil public", "GET portfolio"]),
    ]
    for left, title, lines in boxes:
        add_uml_box(slide, left, Inches(1.6), Inches(2.9), Inches(2.4), title, lines)
    add_uml_box(slide, Inches(0.4), Inches(4.5), Inches(12.5), Inches(1.5), "/api/admin (extrait)", [
        "KYC review · users · providers featured · spotlight · stats · export · subscription override",
        "Paiement abonnement simulé (SUBSCRIPTION_PAYMENT_SIMULATE) — Orange / MVola / Airtel",
    ], LIGHT)
    add_footer(slide)


def slide_features(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, SECONDARY)
    add_header_bar(slide, "Fonctionnalités implémentées")
    done = [
        "Marketplace : annonces services + demandes clients + réponses prestataires",
        "Réservations avec statuts (PENDING → COMPLETED) et avis",
        "Messagerie + offres de prix négociables dans les conversations",
        "KYC prestataire (CIN recto/verso) + validation admin",
        "Portfolio prestataire avec commentaires",
        "Abonnements prestataire (simulation paiement Mobile Money)",
        "Notifications in-app, email et push navigateur",
        "Photos de couverture annonces (1 image max, stockage local)",
        "SEO avancé : sitemap, robots, JSON-LD, pages catégories",
        "Mode sombre · Navigation dashboard avec état actif",
        "Espace admin : modération, spotlight, export, statistiques",
    ]
    add_bullets(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(5.5), "État fonctionnel", done, title_size=24, bullet_size=14)
    add_footer(slide)


def slide_brand(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, SECONDARY)
    add_header_bar(slide, "Charte graphique")
    colors = [
        (BRAND_PRIMARY, "Primaire — Bleu canard", "Boutons, liens, navbar, badges (#069494)"),
        (BRAND_SECONDARY, "Secondaire — Blanc", "Fonds clairs, texte sur primaire (#ffffff)"),
        (BRAND_TERTIARY, "Tertiaire — Pêche", "Accents chaleureux, notifications (#FAD0C9)"),
    ]
    x = Inches(0.8)
    for hex_val, name, desc in colors:
        add_rect(slide, x, Inches(2.0), Inches(3.5), Inches(2.5), hex_rgb(hex_val), PRIMARY if hex_val != BRAND_SECONDARY else MUTED_COLOR)
        add_text_box(slide, x, Inches(4.7), Inches(3.5), Inches(0.4), name, size=16, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
        add_text_box(slide, x, Inches(5.1), Inches(3.5), Inches(0.5), hex_val, size=14, color=INK_COLOR, align=PP_ALIGN.CENTER)
        add_text_box(slide, x, Inches(5.5), Inches(3.5), Inches(0.8), desc, size=11, color=MUTED_COLOR, align=PP_ALIGN.CENTER)
        x += Inches(4.1)
    add_bullets(
        slide, Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.9),
        "",
        ["Tailwind v4 · palette brand-* et tertiary-* dans globals.css · shadcn/ui · Lucide icons"],
        title_size=1, bullet_size=11,
    )
    add_footer(slide)


def slide_status(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT)
    add_header_bar(slide, "État actuel du projet")
    metrics = [
        ("60", "Routes API"),
        ("18", "Modèles Prisma"),
        ("3", "Rôles utilisateur"),
        ("E2E", "Tests Playwright"),
    ]
    x = Inches(0.6)
    for val, label in metrics:
        add_rect(slide, x, Inches(1.5), Inches(2.8), Inches(1.5), SECONDARY, PRIMARY)
        add_text_box(slide, x, Inches(1.65), Inches(2.8), Inches(0.8), val, size=36, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
        add_text_box(slide, x, Inches(2.5), Inches(2.8), Inches(0.4), label, size=14, color=INK_COLOR, align=PP_ALIGN.CENTER)
        x += Inches(3.1)
    add_bullets(
        slide, Inches(0.6), Inches(3.3), Inches(5.8), Inches(3.5),
        "Prêt pour démo",
        [
            "Parcours complet client & prestataire",
            "Auth JWT + vérification email",
            "Publication et recherche d'annonces",
            "Messagerie et réservations",
            "Dashboard admin opérationnel",
        ],
        bullet_size=14,
    )
    add_bullets(
        slide, Inches(6.8), Inches(3.3), Inches(6), Inches(3.5),
        "Limitations connues",
        [
            "Paiements réels non branchés",
            "Stockage fichiers local uniquement",
            "Pas de temps réel (WebSocket)",
            "Pas de PWA / i18n malgache",
            "CI/CD et README à compléter",
        ],
        bullet_size=14,
    )
    add_footer(slide)


def slide_roadmap(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, SECONDARY)
    add_header_bar(slide, "Feuille de route — À faire")
    phases = [
        ("Priorité haute", TERTIARY, [
            "Intégration paiements Mobile Money réels (Orange, MVola, Airtel)",
            "Alimenter le modèle Transaction depuis les réservations",
            "Stockage cloud (S3 / Supabase Storage) pour médias",
            "Rejet KYC formel côté admin + notification prestataire",
        ]),
        ("Priorité moyenne", LIGHT, [
            "Messagerie temps réel (WebSocket ou SSE)",
            "PWA : manifest, service worker, installation mobile",
            "Internationalisation (français + malgache)",
            "Pipeline CI/CD + couverture tests unitaires",
        ]),
        ("Priorité basse", SECONDARY, [
            "README projet et documentation déploiement",
            "Analytics (PostHog ou équivalent)",
            "Optimisation performances (cache, images CDN)",
            "Refonte logo officielle (proposition existante)",
        ]),
    ]
    x = Inches(0.5)
    for title, bg, items in phases:
        add_rect(slide, x, Inches(1.4), Inches(4.0), Inches(5.3), bg, PRIMARY)
        add_text_box(slide, x + Inches(0.15), Inches(1.55), Inches(3.7), Inches(0.45), title, size=18, bold=True, color=PRIMARY)
        y = Inches(2.1)
        for item in items:
            add_text_box(slide, x + Inches(0.15), y, Inches(3.7), Inches(0.95), f"• {item}", size=12, color=INK_COLOR)
            y += Inches(1.05)
        x += Inches(4.2)
    add_footer(slide)


def slide_thanks(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PRIMARY)
    add_rect(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.25), TERTIARY)
    add_text_box(slide, Inches(0.8), Inches(2.5), Inches(11), Inches(1), "Merci", size=48, bold=True, color=SECONDARY, align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.8), Inches(3.8), Inches(11), Inches(0.6), "Tairo ampio — Marketplace de services à Madagascar", size=22, color=SECONDARY, align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.8), Inches(4.8), Inches(11), Inches(0.5), "TairoTairo · Next.js · Prisma · PostgreSQL", size=16, color=TERTIARY, align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(0), Inches(7.0), Inches(13.33), Inches(0.5), TERTIARY)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide_cover(prs)
    slide_vision(prs)
    slide_stack(prs)
    slide_architecture(prs)
    slide_prisma_uml(prs)
    slide_api_overview(prs)
    slide_api_uml_auth(prs)
    slide_api_uml_marketplace(prs)
    slide_api_uml_messaging(prs)
    slide_api_uml_provider(prs)
    slide_features(prs)
    slide_brand(prs)
    slide_status(prs)
    slide_roadmap(prs)
    slide_thanks(prs)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"Présentation générée : {OUTPUT}")
    print(f"Slides : {len(prs.slides)}")


if __name__ == "__main__":
    main()
